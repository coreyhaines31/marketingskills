"""
Generate Envision x UKHS pitch deck as a PowerPoint file.
Output: envision-ukhs-pitch.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Envision brand colors
NAVY = RGBColor(0x00, 0x30, 0x87)
GREEN = RGBColor(0x78, 0xBE, 0x21)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF8, 0xFA, 0xFC)
SLATE = RGBColor(0x47, 0x56, 0x69)
DARK_NAVY = RGBColor(0x00, 0x1D, 0x54)
# UKHS brand colors
UKHS_BLUE = RGBColor(0x00, 0x33, 0x66)
UKHS_RED = RGBColor(0xE8, 0x00, 0x0D)

W = Inches(13.33)
H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def fill_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, left, top, width, height, fill=None, line=None, line_width=Pt(1)):
    s = slide.shapes.add_shape(1, left, top, width, height)
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
        s.line.width = line_width
    else:
        s.line.fill.background()
    return s


def txt(slide, left, top, width, height, text,
        size=12, bold=False, color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return box


# ─────────────────────────────────────────────
# SLIDE 1: Hero
# ─────────────────────────────────────────────
def slide_hero(prs):
    slide = blank_slide(prs)
    fill_bg(slide, NAVY)
    rect(slide, 0, 0, W, Inches(0.08), GREEN)

    txt(slide, Inches(1), Inches(1.2), Inches(11.33), Inches(0.4),
        "EMPOWERING THE UNIVERSITY OF KANSAS HEALTH SYSTEM",
        size=8, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    txt(slide, Inches(1), Inches(1.85), Inches(11.33), Inches(2.2),
        "PRECISION IN\nCARE DELIVERY.",
        size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    txt(slide, Inches(2), Inches(4.1), Inches(9.33), Inches(1.2),
        "Integrating the healthcare print lifecycle to ensure brand integrity,\n"
        "patient data security, and multi-million dollar operational savings.",
        size=15, color=RGBColor(0xCC, 0xD5, 0xE0), align=PP_ALIGN.CENTER)

    txt(slide, Inches(0.5), H - Inches(0.65), Inches(5), Inches(0.35),
        "ENVISION  PRINT SOLUTIONS  |  UKHS STRATEGIC PARTNERSHIP",
        size=7, bold=True, color=RGBColor(0x80, 0x9A, 0xBF))
    rect(slide, 0, H - Inches(0.08), W, Inches(0.08), GREEN)


# ─────────────────────────────────────────────
# SLIDE 2: Clinical Challenges
# ─────────────────────────────────────────────
def slide_challenges(prs):
    slide = blank_slide(prs)
    fill_bg(slide, WHITE)
    rect(slide, 0, 0, W, Inches(0.08), GREEN)

    txt(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.35),
        "CLINICAL LANDSCAPE", size=9, bold=True, color=GREEN)
    txt(slide, Inches(0.8), Inches(0.85), Inches(6.5), Inches(1.3),
        "The Print Challenges of\nModern Health Systems.",
        size=32, bold=True, color=NAVY)
    txt(slide, Inches(7.8), Inches(1.1), Inches(5), Inches(0.9),
        "Managing brand consistency across an academic medical center\n"
        "and hundreds of clinics requires high-performance logistics.",
        size=11, color=SLATE)

    challenges = [
        ("01", "Fragmented Clinic Support",
         "Disparate vendors across regional clinics leading to brand drift and hidden logistics costs."),
        ("02", "Data Privacy Standards",
         "Ensuring 100% HIPAA and PHI compliance for every personalized member communication."),
        ("03", "Manual Patient Comm",
         "Slow workflows for critical patient education and discharge materials leading to care gaps."),
        ("04", "Dynamic Localization",
         "Managing facility-specific materials for hundreds of doctors and dozens of specialties."),
        ("05", "Inventory Risk",
         "Obsolete brochures and forms taking up valuable clinical floor space."),
        ("06", "Lack of Optimization",
         "No system-wide analytics to identify cost-saving opportunities and waste reduction."),
    ]

    col_w = Inches(4.1)
    row_h = Inches(2.3)
    sx, sy = Inches(0.4), Inches(2.3)
    for i, (num, title, desc) in enumerate(challenges):
        col, row = i % 3, i // 3
        x = sx + col * (col_w + Inches(0.2))
        y = sy + row * (row_h + Inches(0.12))
        rect(slide, x, y, col_w, row_h, LIGHT_GRAY)
        txt(slide, x + Inches(0.25), y + Inches(0.2), Inches(0.7), Inches(0.55),
            num, size=24, bold=True, color=GREEN)
        txt(slide, x + Inches(0.25), y + Inches(0.75), col_w - Inches(0.5), Inches(0.45),
            title, size=13, bold=True, color=NAVY)
        txt(slide, x + Inches(0.25), y + Inches(1.2), col_w - Inches(0.5), Inches(0.95),
            desc, size=10, color=SLATE)


# ─────────────────────────────────────────────
# SLIDE 3: Clinical Methodology
# ─────────────────────────────────────────────
def slide_methodology(prs):
    slide = blank_slide(prs)
    fill_bg(slide, DARK_NAVY)
    rect(slide, 0, 0, W, Inches(0.08), GREEN)

    txt(slide, Inches(0.8), Inches(0.7), Inches(5), Inches(0.35),
        "OUR METHODOLOGY", size=9, bold=True, color=GREEN)
    txt(slide, Inches(0.8), Inches(1.1), Inches(5.5), Inches(1.5),
        "Integrating the Clinical\nPrint Lifecycle.",
        size=34, bold=True, color=WHITE)

    steps = [
        ("1", "Discovery & Compliance Strategy",
         "Detailed assessment of the KU Health System supply chain, prioritizing HIPAA protocols and brand cohesion."),
        ("2", "Validated Production",
         "G7 certified technology with 100% inspection to ensure every patient touchpoint meets the system's standards."),
        ("3", "Digital Portal Architecture",
         "Web-to-print portals that allow clinics to order localized, pre-approved materials instantly, reducing overhead."),
    ]
    for i, (num, title, desc) in enumerate(steps):
        y = Inches(2.7) + i * Inches(1.45)
        circ = slide.shapes.add_shape(9, Inches(0.8), y, Inches(0.55), Inches(0.55))
        circ.fill.background()
        circ.line.color.rgb = GREEN
        circ.line.width = Pt(1.5)
        txt(slide, Inches(0.88), y + Inches(0.06), Inches(0.4), Inches(0.45),
            num, size=13, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
        txt(slide, Inches(1.6), y, Inches(4.8), Inches(0.42),
            title, size=14, bold=True, color=WHITE)
        txt(slide, Inches(1.6), y + Inches(0.42), Inches(4.8), Inches(0.65),
            desc, size=10, color=RGBColor(0x94, 0xA3, 0xB8))

    # Right panel
    rect(slide, Inches(7.2), Inches(1.5), Inches(5.5), Inches(4.5), DARK_NAVY)
    box = rect(slide, Inches(7.6), Inches(2.0), Inches(4.7), Inches(3.5),
               line=GREEN, line_width=Pt(3))
    txt(slide, Inches(7.8), Inches(2.4), Inches(4.3), Inches(0.9),
        "SECURITY.", size=48, bold=True,
        color=RGBColor(0x40, 0x55, 0x7A), align=PP_ALIGN.CENTER)
    txt(slide, Inches(7.8), Inches(3.1), Inches(4.3), Inches(0.9),
        "INTEGRITY.", size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(slide, Inches(7.8), Inches(4.1), Inches(4.3), Inches(0.5),
        "ENGINEERED FOR UKHS", size=9, bold=True, color=GREEN, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────
# SLIDE 4: Operational Impact Metrics
# ─────────────────────────────────────────────
def slide_results(prs):
    slide = blank_slide(prs)
    fill_bg(slide, WHITE)
    rect(slide, 0, 0, W, Inches(0.08), GREEN)

    txt(slide, 0, Inches(0.8), W, Inches(0.8),
        "Operational Impact for UKHS.",
        size=38, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    rect(slide, Inches(5.42), Inches(1.65), Inches(2.5), Inches(0.1), GREEN)

    metrics = [
        ("25-40%", "SYSTEM-WIDE SAVINGS"),
        ("100%",   "PHI/HIPAA COMPLIANCE"),
        ("20%+",   "PATIENT ENGAGEMENT"),
        ("60%",    "INVENTORY WASTE REDUCTION"),
    ]
    col_w = Inches(3.0)
    sx = Inches(0.42)
    for i, (value, label) in enumerate(metrics):
        x = sx + i * (col_w + Inches(0.11))
        rect(slide, x, Inches(2.2), col_w, Inches(4.5), LIGHT_GRAY)
        txt(slide, x, Inches(3.3), col_w, Inches(1.4),
            value, size=50, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        txt(slide, x, Inches(4.8), col_w, Inches(0.65),
            label, size=9, bold=True, color=SLATE, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────
# SLIDE 5: Testimonial
# ─────────────────────────────────────────────
def slide_testimonial(prs):
    slide = blank_slide(prs)
    fill_bg(slide, LIGHT_GRAY)
    rect(slide, 0, 0, W, Inches(0.08), GREEN)

    txt(slide, Inches(0.8), Inches(0.9), Inches(0.8), Inches(0.9),
        "\u201c", size=72, bold=True, color=GREEN)
    txt(slide, Inches(0.8), Inches(1.6), Inches(6.2), Inches(3.2),
        "Envision has unified our brand voice across 300+ care locations, "
        "delivering significant annual savings while ensuring our patient data "
        "remains absolutely secure.",
        size=21, bold=True, color=NAVY)

    rect(slide, Inches(0.8), Inches(5.05), Inches(1.0), Inches(0.08), GREEN)
    txt(slide, Inches(2.1), Inches(4.85), Inches(5), Inches(0.45),
        "CLINICAL OPERATIONS LEAD", size=13, bold=True, color=NAVY)
    txt(slide, Inches(2.1), Inches(5.3), Inches(5), Inches(0.4),
        "Top-Tier Academic Medical Center", size=11, color=SLATE)

    rect(slide, Inches(7.5), 0, Inches(5.83), H, NAVY)
    txt(slide, Inches(7.8), Inches(2.2), Inches(5.0), Inches(0.8),
        "300+ CARE LOCATIONS", size=10, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    txt(slide, Inches(7.8), Inches(3.0), Inches(5.0), Inches(1.5),
        "One unified brand voice.\nComplete PHI security.\nProven multi-million dollar savings.",
        size=14, color=RGBColor(0xCC, 0xD5, 0xE0), align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────
# SLIDES 6-8: Case Studies
# ─────────────────────────────────────────────
def slide_case_study(prs, title, challenge, solution, results):
    slide = blank_slide(prs)
    fill_bg(slide, WHITE)
    rect(slide, 0, 0, W, Inches(0.08), GREEN)

    rect(slide, 0, 0, Inches(5.5), H, LIGHT_GRAY)
    txt(slide, Inches(0.6), Inches(0.7), Inches(4.5), Inches(0.35),
        "PARTNER SUCCESS STORY", size=8, bold=True, color=GREEN)
    txt(slide, Inches(0.6), Inches(1.1), Inches(4.5), Inches(1.1),
        title, size=28, bold=True, color=NAVY)
    txt(slide, Inches(0.6), Inches(2.25), Inches(4.5), Inches(1.5),
        challenge, size=12, color=SLATE)

    rect(slide, Inches(0.6), Inches(3.9), Inches(4.5), Inches(1.85), WHITE)
    rect(slide, Inches(0.6), Inches(3.9), Inches(0.07), Inches(1.85), GREEN)
    txt(slide, Inches(0.85), Inches(4.05), Inches(4), Inches(0.35),
        "ENVISION'S SOLUTION", size=8, bold=True, color=RGBColor(0x94, 0xA3, 0xB8))
    txt(slide, Inches(0.85), Inches(4.45), Inches(4), Inches(1.0),
        solution, size=12, bold=True, color=NAVY)

    rect(slide, Inches(5.5), 0, Inches(7.83), H, DARK_NAVY)
    stat_w = Inches(2.3)
    for i, (stat, label) in enumerate(results):
        sx = Inches(5.8) + i * (stat_w + Inches(0.15))
        sy = Inches(2.4)
        rect(slide, sx, sy, stat_w, Inches(2.2), RGBColor(0x00, 0x28, 0x6A))
        rect(slide, sx, sy, stat_w, Inches(0.07), GREEN)
        txt(slide, sx, sy + Inches(0.5), stat_w, Inches(0.8),
            stat, size=30, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(slide, sx, sy + Inches(1.35), stat_w, Inches(0.5),
            label, size=8, bold=True,
            color=RGBColor(0xCC, 0xD5, 0xE0), align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────
# SLIDE 9: CTA
# ─────────────────────────────────────────────
def slide_cta(prs):
    slide = blank_slide(prs)
    fill_bg(slide, WHITE)
    rect(slide, 0, 0, W, Inches(0.08), GREEN)

    txt(slide, Inches(1), Inches(0.85), Inches(11.33), Inches(1.7),
        "Advancing the power of\nmedicine through print.",
        size=36, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    txt(slide, Inches(2), Inches(2.6), Inches(9.33), Inches(0.85),
        "Let's schedule a strategic audit for The University of Kansas Health System\n"
        "to streamline clinical operations and boost brand consistency.",
        size=13, color=SLATE, align=PP_ALIGN.CENTER)

    items = ["Standardize System Brand", "Automate Clinic Ordering"]
    col_w = Inches(4.5)
    for i, item in enumerate(items):
        x = Inches(1.8) + i * (col_w + Inches(0.4))
        rect(slide, x, Inches(3.65), col_w, Inches(0.75), LIGHT_GRAY)
        rect(slide, x + Inches(0.15), Inches(3.78), Inches(0.45), Inches(0.45), GREEN)
        txt(slide, x + Inches(0.15), Inches(3.78), Inches(0.45), Inches(0.45),
            "\u2713", size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txt(slide, x + Inches(0.75), Inches(3.77), col_w - Inches(0.9), Inches(0.5),
            item, size=13, bold=True, color=NAVY)

    rect(slide, Inches(4.2), Inches(5.05), Inches(4.93), Inches(0.85), NAVY)
    txt(slide, Inches(4.2), Inches(5.15), Inches(4.93), Inches(0.65),
        "REQUEST ASSESSMENT",
        size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────
# SLIDE 10: Contact
# ─────────────────────────────────────────────
def slide_contact(prs):
    slide = blank_slide(prs)
    fill_bg(slide, DARK_NAVY)
    rect(slide, 0, 0, W, Inches(0.08), GREEN)

    txt(slide, Inches(0.8), Inches(0.7), Inches(4), Inches(0.35),
        "LOCAL COMMITMENT", size=9, bold=True, color=GREEN)
    txt(slide, Inches(0.8), Inches(1.1), Inches(3.5), Inches(1.6),
        "Support for\nKansas Leaders.",
        size=32, bold=True, color=WHITE)
    txt(slide, Inches(0.8), Inches(2.85), Inches(3.3), Inches(1.2),
        "Our leadership team is dedicated to supporting the unique mission of "
        "The University of Kansas Health System across the region.",
        size=11, color=RGBColor(0x94, 0xA3, 0xB8))

    contacts = [
        ("Sebastian Zahr", "Chief Revenue & Strategy", "316-655-4243",
         "Sebastian.Zahr@Envisionus.com"),
        ("Kathy Vines", "VP Strategic Partnerships", "316-425-7105",
         "Kathy.Vines@Envisionus.com"),
        ("Rob Narron", "Strategic Business Dev", "316-425-7268",
         "Rob.Narron@Envisionus.com"),
    ]
    col_w = Inches(2.8)
    sx = Inches(4.4)
    for i, (name, title, phone, email) in enumerate(contacts):
        x = sx + i * (col_w + Inches(0.35))
        rect(slide, x, Inches(1.2), Inches(0.04), Inches(5.5), GREEN)
        txt(slide, x + Inches(0.2), Inches(1.3), col_w, Inches(0.5),
            name, size=16, bold=True, color=WHITE)
        txt(slide, x + Inches(0.2), Inches(1.85), col_w, Inches(0.5),
            title.upper(), size=7.5, bold=True, color=GREEN)
        txt(slide, x + Inches(0.2), Inches(2.9), col_w, Inches(0.4),
            phone, size=11, color=RGBColor(0x94, 0xA3, 0xB8))
        txt(slide, x + Inches(0.2), Inches(3.35), col_w, Inches(0.4),
            email, size=9, color=RGBColor(0x7A, 0x90, 0xAA))


# ─────────────────────────────────────────────
# Build the deck
# ─────────────────────────────────────────────
def main():
    prs = new_prs()

    slide_hero(prs)
    slide_challenges(prs)
    slide_methodology(prs)
    slide_results(prs)
    slide_testimonial(prs)

    slide_case_study(
        prs,
        "Regional Academic Center",
        "A multi-hospital academic system struggled with brand dilution across 45 specialty "
        "clinics and rising PHI compliance risks.",
        "Consolidated on-demand ordering with dynamic clinic localization and audited HIPAA workflows.",
        [("$3.2M", "ANNUAL SAVINGS"), ("22%", "PATIENT RETURN"), ("100%", "AUDIT PASS")]
    )
    slide_case_study(
        prs,
        "Cancer Center Network",
        "Highly complex, personalized education packets were taking weeks to print and assemble, "
        "delaying care initiation.",
        "Triggered on-demand digital production with automated kitting for personalized patient care plans.",
        [("80%", "FASTER CYCLE"), ("15%", "PATIENT SAT LIFT"), ("ZERO", "WASTE")]
    )
    slide_case_study(
        prs,
        "Health Plan Services",
        "Member communication for Open Enrollment lacked dynamic personalization, resulting in "
        "low plan retention.",
        "Full-color variable data direct mail utilizing predictive modeling to offer personalized "
        "benefits highlights.",
        [("31%", "HIGHER RETENTION"), ("14x", "RETURN ON AD SPEND"), ("360°", "VIEW ENABLED")]
    )

    slide_cta(prs)
    slide_contact(prs)

    out = "envision-ukhs-pitch.pptx"
    prs.save(out)
    print(f"Saved: {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
