"""
Generate Envision Print Solutions pitch deck as a PowerPoint file.
Output: envision-print-pitch.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Brand colors
NAVY = RGBColor(0x00, 0x30, 0x87)
GREEN = RGBColor(0x78, 0xBE, 0x21)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF8, 0xFA, 0xFC)
SLATE = RGBColor(0x47, 0x56, 0x69)
DARK_NAVY = RGBColor(0x00, 0x1D, 0x54)

W = Inches(13.33)
H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]  # Completely blank
    return prs.slides.add_slide(blank_layout)


def fill_slide_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False,
                color=WHITE, align=PP_ALIGN.LEFT, font_name="Calibri", wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_label(slide, left, top, width, text, color=GREEN, size=9):
    add_textbox(slide, left, top, width, Inches(0.35), text,
                font_size=size, bold=True, color=color, align=PP_ALIGN.LEFT)


# ─────────────────────────────────────────────
# SLIDE 1: Hero
# ─────────────────────────────────────────────
def slide_hero(prs):
    slide = blank_slide(prs)
    fill_slide_bg(slide, NAVY)

    # Green accent line at top
    add_rect(slide, 0, 0, W, Inches(0.08), GREEN)

    # Eyebrow label
    add_textbox(slide, Inches(1), Inches(1.4), Inches(11.33), Inches(0.4),
                "PARTNERING FOR PERFORMANCE",
                font_size=9, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    # Main title
    add_textbox(slide, Inches(1), Inches(2.0), Inches(11.33), Inches(2.4),
                "HIGH-IMPACT PRINTING.",
                font_size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Subtitle
    add_textbox(slide, Inches(2), Inches(4.5), Inches(9.33), Inches(1.2),
                "Integrating every stage of the print lifecycle to drive quality,\n"
                "efficiency, and multi-million dollar business results.",
                font_size=16, bold=False, color=RGBColor(0xCC, 0xD5, 0xE0),
                align=PP_ALIGN.CENTER)

    # Bottom green bar
    add_rect(slide, 0, H - Inches(0.08), W, Inches(0.08), GREEN)

    # Logo text area bottom left
    add_textbox(slide, Inches(0.5), H - Inches(0.7), Inches(3), Inches(0.4),
                "ENVISION  PRINT SOLUTIONS",
                font_size=8, bold=True, color=RGBColor(0x80, 0x9A, 0xBF),
                align=PP_ALIGN.LEFT)


# ─────────────────────────────────────────────
# SLIDE 2: Market Landscape — Challenges
# ─────────────────────────────────────────────
def slide_challenges(prs):
    slide = blank_slide(prs)
    fill_slide_bg(slide, WHITE)
    add_rect(slide, 0, 0, W, Inches(0.08), GREEN)

    add_label(slide, Inches(0.8), Inches(0.5), Inches(4), "MARKET LANDSCAPE")

    add_textbox(slide, Inches(0.8), Inches(0.85), Inches(6), Inches(1.2),
                "The Print Challenges\nFacing Enterprises.",
                font_size=32, bold=True, color=NAVY)

    add_textbox(slide, Inches(7.5), Inches(1.1), Inches(5.3), Inches(0.9),
                "Fragmented supply chains and manual processes are\nsilent profit killers in the modern enterprise.",
                font_size=12, color=SLATE)

    challenges = [
        ("01", "Fragmented Supply Chain",
         "No single point of accountability leads to hidden costs and vendor fatigue."),
        ("02", "Quality Inconsistency",
         "Unreliable color accuracy erodes brand trust across different markets."),
        ("03", "Manual Bottlenecks",
         "Archaic workflows and ordering systems create errors and project delays."),
        ("04", "Scale Barriers",
         "Inability to execute 1-to-1 personalized campaigns at enterprise volumes."),
        ("05", "Inventory Waste",
         "Capital locked in warehousing obsolete assets that end up in landfills."),
        ("06", "Reporting Blind Spots",
         "Lack of analytics to optimize ROI and justify print marketing spend."),
    ]

    cols = 3
    col_w = Inches(4.1)
    row_h = Inches(2.3)
    start_x = Inches(0.4)
    start_y = Inches(2.3)

    for i, (num, title, desc) in enumerate(challenges):
        col = i % cols
        row = i // cols
        x = start_x + col * (col_w + Inches(0.2))
        y = start_y + row * (row_h + Inches(0.12))

        add_rect(slide, x, y, col_w, row_h, LIGHT_GRAY)
        add_textbox(slide, x + Inches(0.25), y + Inches(0.2), Inches(0.7), Inches(0.6),
                    num, font_size=24, bold=True, color=GREEN)
        add_textbox(slide, x + Inches(0.25), y + Inches(0.75), col_w - Inches(0.5), Inches(0.45),
                    title, font_size=13, bold=True, color=NAVY)
        add_textbox(slide, x + Inches(0.25), y + Inches(1.2), col_w - Inches(0.5), Inches(0.95),
                    desc, font_size=10, color=SLATE)


# ─────────────────────────────────────────────
# SLIDE 3: Methodology
# ─────────────────────────────────────────────
def slide_methodology(prs):
    slide = blank_slide(prs)
    fill_slide_bg(slide, DARK_NAVY)
    add_rect(slide, 0, 0, W, Inches(0.08), GREEN)

    # Left panel
    add_label(slide, Inches(0.8), Inches(0.7), Inches(5), "OUR METHODOLOGY", color=GREEN)

    add_textbox(slide, Inches(0.8), Inches(1.1), Inches(5.5), Inches(1.4),
                "Integrating the\nPrint Lifecycle.",
                font_size=36, bold=True, color=WHITE)

    steps = [
        ("1", "Discovery & Strategy",
         "Cost reduction audits and goal alignment to ensure every piece has a purpose."),
        ("2", "Automated Production",
         "G7 certified tech paired with ISO 9001 quality control for flawless execution."),
        ("3", "Technology & Analytics",
         "Web-to-print portals and attribution modeling to measure true campaign ROI."),
    ]

    for i, (num, title, desc) in enumerate(steps):
        y = Inches(2.7) + i * Inches(1.45)
        # Circle number
        circ = slide.shapes.add_shape(9, Inches(0.8), y, Inches(0.55), Inches(0.55))  # oval
        circ.fill.background()
        circ.line.color.rgb = GREEN
        circ.line.width = Pt(1.5)
        add_textbox(slide, Inches(0.88), y + Inches(0.05), Inches(0.4), Inches(0.45),
                    num, font_size=13, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

        add_textbox(slide, Inches(1.6), y, Inches(4.5), Inches(0.4),
                    title, font_size=14, bold=True, color=WHITE)
        add_textbox(slide, Inches(1.6), y + Inches(0.4), Inches(4.5), Inches(0.6),
                    desc, font_size=10, color=RGBColor(0x94, 0xA3, 0xB8))

    # Right panel highlight box
    add_rect(slide, Inches(7.2), Inches(1.5), Inches(5.5), Inches(4.5), DARK_NAVY)
    box = add_rect(slide, Inches(7.6), Inches(2.0), Inches(4.7), Inches(3.5), None, GREEN)
    box.line.width = Pt(3)

    add_textbox(slide, Inches(7.8), Inches(2.4), Inches(4.3), Inches(0.9),
                "SPEED.", font_size=48, bold=True,
                color=RGBColor(0x40, 0x55, 0x7A), align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(7.8), Inches(3.1), Inches(4.3), Inches(0.9),
                "PRECISION.", font_size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(7.8), Inches(4.1), Inches(4.3), Inches(0.5),
                "ENGINEERED FOR ENTERPRISE",
                font_size=9, bold=True, color=GREEN, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────
# SLIDE 4: Measurable Outcomes
# ─────────────────────────────────────────────
def slide_results(prs):
    slide = blank_slide(prs)
    fill_slide_bg(slide, WHITE)
    add_rect(slide, 0, 0, W, Inches(0.08), GREEN)

    add_textbox(slide, 0, Inches(0.8), W, Inches(0.8),
                "Measurable Outcomes.",
                font_size=40, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # Green underline
    add_rect(slide, Inches(5.92), Inches(1.65), Inches(1.5), Inches(0.1), GREEN)

    metrics = [
        ("25-40%", "SPEND REDUCTION"),
        ("50%+", "SPEED TO MARKET"),
        ("20%+", "RESPONSE BOOST"),
        ("100%", "ON-TIME DELIVERY"),
    ]

    col_w = Inches(3.0)
    start_x = Inches(0.42)
    y_top = Inches(2.2)

    for i, (value, label) in enumerate(metrics):
        x = start_x + i * (col_w + Inches(0.11))
        add_rect(slide, x, y_top, col_w, Inches(4.5), LIGHT_GRAY)
        add_textbox(slide, x, y_top + Inches(1.2), col_w, Inches(1.5),
                    value, font_size=52, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        add_textbox(slide, x, y_top + Inches(2.8), col_w, Inches(0.6),
                    label, font_size=10, bold=True, color=SLATE, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────
# SLIDE 5: Testimonial
# ─────────────────────────────────────────────
def slide_testimonial(prs):
    slide = blank_slide(prs)
    fill_slide_bg(slide, LIGHT_GRAY)
    add_rect(slide, 0, 0, W, Inches(0.08), GREEN)

    # Left quote panel
    add_textbox(slide, Inches(0.8), Inches(0.9), Inches(0.8), Inches(0.9),
                "\u201c", font_size=72, bold=True, color=GREEN)

    add_textbox(slide, Inches(0.8), Inches(1.6), Inches(6.2), Inches(3.2),
                "Envision has transformed our print supply chain into a competitive "
                "advantage and profit center. Their solutions have delivered multi-million "
                "dollar savings.",
                font_size=22, bold=True, color=NAVY)

    # Green divider line
    add_rect(slide, Inches(0.8), Inches(5.1), Inches(1.0), Inches(0.08), GREEN)

    add_textbox(slide, Inches(2.1), Inches(4.9), Inches(5), Inches(0.45),
                "CHIEF MARKETING OFFICER", font_size=13, bold=True, color=NAVY)
    add_textbox(slide, Inches(2.1), Inches(5.35), Inches(5), Inches(0.4),
                "Fortune 500 Financial Services", font_size=11, color=SLATE)

    # Right panel — dark navy accent
    add_rect(slide, Inches(7.5), 0, Inches(5.83), H, NAVY)
    add_textbox(slide, Inches(7.8), Inches(2.5), Inches(5.0), Inches(0.9),
                "RESULTS THAT MATTER", font_size=11, bold=True, color=GREEN,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(7.8), Inches(3.2), Inches(5.0), Inches(1.2),
                "Trusted by Fortune 500 companies\nacross financial services, healthcare,\nretail and CPG.",
                font_size=14, color=RGBColor(0xCC, 0xD5, 0xE0), align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────
# SLIDES 6-8: Case Studies
# ─────────────────────────────────────────────
def slide_case_study(prs, title, challenge, solution, results):
    slide = blank_slide(prs)
    fill_slide_bg(slide, WHITE)
    add_rect(slide, 0, 0, W, Inches(0.08), GREEN)

    # Left panel
    add_rect(slide, 0, 0, Inches(5.5), H, LIGHT_GRAY)
    add_label(slide, Inches(0.6), Inches(0.7), Inches(4), "SUCCESS STORY")

    add_textbox(slide, Inches(0.6), Inches(1.15), Inches(4.5), Inches(1.0),
                title, font_size=30, bold=True, color=NAVY)

    add_textbox(slide, Inches(0.6), Inches(2.2), Inches(4.5), Inches(1.5),
                challenge, font_size=12, color=SLATE)

    # Solution box
    add_rect(slide, Inches(0.6), Inches(3.9), Inches(4.5), Inches(1.8), WHITE)
    add_textbox(slide, Inches(0.85), Inches(4.05), Inches(4), Inches(0.35),
                "THE SOLUTION", font_size=8, bold=True, color=RGBColor(0x94, 0xA3, 0xB8))
    add_textbox(slide, Inches(0.85), Inches(4.45), Inches(4), Inches(0.95),
                solution, font_size=13, bold=True, color=NAVY)

    # Left border accent on solution box
    add_rect(slide, Inches(0.6), Inches(3.9), Inches(0.07), Inches(1.8), GREEN)

    # Right panel — dark with result stats
    add_rect(slide, Inches(5.5), 0, Inches(7.83), H, DARK_NAVY)

    stat_w = Inches(2.3)
    stat_h = Inches(2.2)
    stat_y = Inches(2.4)
    for i, (stat, label) in enumerate(results):
        sx = Inches(5.8) + i * (stat_w + Inches(0.15))
        add_rect(slide, sx, stat_y, stat_w, stat_h, RGBColor(0x00, 0x28, 0x6A))
        # Green top border
        add_rect(slide, sx, stat_y, stat_w, Inches(0.07), GREEN)
        add_textbox(slide, sx, stat_y + Inches(0.5), stat_w, Inches(0.8),
                    stat, font_size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, sx, stat_y + Inches(1.3), stat_w, Inches(0.5),
                    label, font_size=8, bold=True,
                    color=RGBColor(0xCC, 0xD5, 0xE0), align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────
# SLIDE 9: CTA
# ─────────────────────────────────────────────
def slide_cta(prs):
    slide = blank_slide(prs)
    fill_slide_bg(slide, WHITE)
    add_rect(slide, 0, 0, W, Inches(0.08), GREEN)

    add_textbox(slide, Inches(1.5), Inches(0.8), Inches(10.33), Inches(1.8),
                "Ready to realize your print potential?",
                font_size=36, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(2), Inches(2.5), Inches(9.33), Inches(0.8),
                "Schedule a strategic print assessment to uncover hidden efficiencies\n"
                "and boost your marketing ROI with a partner accountable to your goals.",
                font_size=13, color=SLATE, align=PP_ALIGN.CENTER)

    checklist = [
        "Streamline Fragmented Operations",
        "Ensure Brand Consistency",
    ]
    col_w = Inches(4.5)
    for i, item in enumerate(checklist):
        x = Inches(1.8) + i * (col_w + Inches(0.4))
        add_rect(slide, x, Inches(3.6), col_w, Inches(0.75), LIGHT_GRAY)
        # Green checkbox
        add_rect(slide, x + Inches(0.15), Inches(3.73), Inches(0.45), Inches(0.45), GREEN)
        add_textbox(slide, x + Inches(0.15), Inches(3.73), Inches(0.45), Inches(0.45),
                    "\u2713", font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.75), Inches(3.72), col_w - Inches(0.9), Inches(0.5),
                    item, font_size=13, bold=True, color=NAVY)

    # CTA Button
    btn = add_rect(slide, Inches(4.4), Inches(5.0), Inches(4.53), Inches(0.85), NAVY)
    add_textbox(slide, Inches(4.4), Inches(5.1), Inches(4.53), Inches(0.65),
                "REQUEST ASSESSMENT",
                font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────
# SLIDE 10: Contact
# ─────────────────────────────────────────────
def slide_contact(prs):
    slide = blank_slide(prs)
    fill_slide_bg(slide, DARK_NAVY)
    add_rect(slide, 0, 0, W, Inches(0.08), GREEN)

    add_label(slide, Inches(0.8), Inches(0.7), Inches(4), "LET'S GET STARTED")

    add_textbox(slide, Inches(0.8), Inches(1.1), Inches(3.5), Inches(1.6),
                "Expert Advice.\nProven Results.",
                font_size=32, bold=True, color=WHITE)

    add_textbox(slide, Inches(0.8), Inches(2.8), Inches(3.3), Inches(1.2),
                "Our leadership team is ready to discuss your unique print challenges "
                "and design a custom solution architecture.",
                font_size=11, color=RGBColor(0x94, 0xA3, 0xB8))

    contacts = [
        ("Sebastian Zahr", "Chief Revenue & Strategy", "316-655-4243",
         "Sebastian.Zahr@Envisionus.com"),
        ("Kathy Vines", "VP Strategic Partnerships", "316-425-7105",
         "Kathy.Vines@Envisionus.com"),
        ("Rob Narron", "Strategic Business Dev", "316-425-7268",
         "Rob.Narron@Envisionus.com"),
    ]

    col_w = Inches(2.8)
    start_x = Inches(4.4)
    for i, (name, title, phone, email) in enumerate(contacts):
        x = start_x + i * (col_w + Inches(0.35))
        # Vertical green accent line
        add_rect(slide, x, Inches(1.2), Inches(0.04), Inches(5.5), GREEN)
        add_textbox(slide, x + Inches(0.2), Inches(1.3), col_w, Inches(0.5),
                    name, font_size=16, bold=True, color=WHITE)
        add_textbox(slide, x + Inches(0.2), Inches(1.85), col_w, Inches(0.5),
                    title.upper(), font_size=7.5, bold=True, color=GREEN)
        add_textbox(slide, x + Inches(0.2), Inches(2.9), col_w, Inches(0.4),
                    phone, font_size=11, color=RGBColor(0x94, 0xA3, 0xB8))
        add_textbox(slide, x + Inches(0.2), Inches(3.35), col_w, Inches(0.4),
                    email, font_size=9, color=RGBColor(0x7A, 0x90, 0xAA))


# ─────────────────────────────────────────────
# Build the deck
# ─────────────────────────────────────────────
def main():
    prs = new_prs()

    slide_hero(prs)
    slide_challenges(prs)
    slide_methodology(prs)
    slide_results(prs)
    slide_testimonial(prs)

    # Three case studies
    slide_case_study(
        prs,
        "Leading Healthcare",
        "Fragmented operations and manual processes were hindering HIPAA-compliant "
        "member communications at scale.",
        "Centralized on-demand dynamic printing with closed-loop attribution.",
        [("$3M", "SAVINGS"), ("25%", "RESPONSE LIFT"), ("100%", "PHI SAFE")]
    )
    slide_case_study(
        prs,
        "Global CPG",
        "Slow time-to-market for new SKUs and high fixed costs for local market "
        "plate/proofing.",
        "Digital-first network and automated asset versioning for instant localization.",
        [("60%", "SPEED INCREASE"), ("43%", "COST REDUCTION"), ("12%", "SALES LIFT")]
    )
    slide_case_study(
        prs,
        "National Retailer",
        "Wasted spend on mass circulars with no personalization or tracking capabilities.",
        "Variable direct mail with lookalike modeling and unique digital promo integration.",
        [("4.7x", "ROAS"), ("22%", "REDEMPTION RATE"), ("360°", "ATTRIBUTION")]
    )

    slide_cta(prs)
    slide_contact(prs)

    out = "envision-print-pitch.pptx"
    prs.save(out)
    print(f"Saved: {out}  ({prs.slides.__len__()} slides)")


if __name__ == "__main__":
    main()
